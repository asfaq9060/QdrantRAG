# backend/utils/extractors.py
import io
import json
import os
import pandas as pd
from PIL import Image
import PyPDF2
from docx import Document
from pptx import Presentation
from .genai_wrapper import genai_generate_text, transcribe_audio_bytes
import uuid

def extract_text_from_image(image_bytes: bytes) -> str:
    try:
        image = Image.open(io.BytesIO(image_bytes))
        prompt = [
            "Extract ALL visible text from this image exactly as it appears. Include handwriting, labels, numbers. Return ONLY the text.",
            image
        ]
        return genai_generate_text(prompt) or "[No text in image]"
    except Exception as e:
        return f"[Image error: {str(e)}]"

def extract_audio_with_gemini(audio_bytes: bytes, filename: str) -> str:
    """
    Write audio bytes to a temp file → upload to Gemini → return transcript.
    """
    try:
        transcript = transcribe_audio_bytes(audio_bytes, filename)
        return transcript or "[No transcript produced]"
    except Exception as e:
        return f"[Audio error: {str(e)}]"

def extract_text(file_bytes: bytes, filename: str, ext: str) -> str:
    try:
        ext = ext.lower()
        if ext in ["png", "jpg", "jpeg", "webp"]:
            return extract_text_from_image(file_bytes)
        if ext in ["mp3", "wav", "m4a", "ogg"]:
            return extract_audio_with_gemini(file_bytes, filename)
        if ext == "pdf":
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        if ext == "docx":
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if ext in ["pptx", "ppt"]:
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        if ext in ["xlsx", "xls"]:
            return pd.read_excel(io.BytesIO(file_bytes)).to_string()
        if ext == "csv":
            return pd.read_csv(io.BytesIO(file_bytes)).to_string()
        if ext == "txt":
            return file_bytes.decode("utf-8", errors="ignore")
        if ext == "json":
            return json.dumps(json.loads(file_bytes.decode("utf-8", errors="ignore")), indent=2)
        return "[File uploaded]"
    except Exception as e:
        return f"[File processed with warning: {str(e)}]"

def chunk_text(text: str, chunk_words: int = 400, step: int = 350):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i+chunk_words])
        if chunk.strip():
            chunks.append(chunk)
        i += step
    return chunks or ["[Empty]"]
